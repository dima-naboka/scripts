
# coding: utf-8

# In[1]:


# import json
# #print(json.dumps(job,indent=4))

# client = dataiku.api_client()
# current_project = client.get_project(dataiku.get_custom_variables()["projectKey"])
# variables = current_project.get_variables()
# general_settings = client.get_general_settings()


# In[2]:


import dataiku
from pptx import Presentation
import io
io_buffer = io.BytesIO()


# In[3]:


def create_pptx_template():
    #writing to buffer NOT FILE
    prs=Presentation()
    lyt=prs.slide_layouts[1] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    title=slide.shapes.title # assigning a title
    subtitle=slide.placeholders[1] # placeholder for subtitle
    title.text="Hey,This is a Slide! How exciting!" # title
    subtitle.text="Really?" # subtitle

    prs.save(io_buffer)


# In[4]:


#upload Dataset as Stream
#https://doc.dataiku.com/dss/latest/python-api/managed_folders.html?highlight=upload_stream#dataiku.Folder.upload_stream

def upload_file_to_sharepoint(managed_folder_id, output_file_name,file_format):
    output_folder = dataiku.Folder(managed_folder_id)
    io_buffer.seek(0)
    output_folder.upload_stream(output_file_name + file_format, io_buffer)


# In[5]:


managed_folder_id = 'Alu1O252'
output_file_name = 'powerpoint_template'
file_format = '.pptx'

create_pptx_template()
upload_file_to_sharepoint(managed_folder_id,output_file_name,file_format)

